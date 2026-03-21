from pptx import Presentation
from pptx.util import Inches, Pt
import io

def generate_ppt(project_data):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = project_data.get('title', 'Project Presentation')
    subtitle.text = f"Presented by Student\nDomain: {project_data.get('domain', '')}"
    
    # Abstract Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Abstract"
    slide.placeholders[1].text_frame.text = project_data.get('abstract', '')[:1000]
    
    # Literature Survey
    if project_data.get('literature_survey'):
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "2. Literature Survey"
        slide.placeholders[1].text_frame.text = project_data.get('literature_survey', '')[:800]

    # Problem Statement
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "3. Problem Statement"
    slide.placeholders[1].text_frame.text = project_data.get('problem_statement', '')
    
    # Key Features
    if project_data.get('features'):
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "4. Key Features"
        tf = slide.placeholders[1].text_frame
        for feature in project_data.get('features', []):
            p = tf.add_paragraph()
            p.text = feature

    # Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "5. System Architecture"
    slide.placeholders[1].text_frame.text = project_data.get('architecture_description', '')
    
    # Technology Stack
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "6. Technology Stack"
    tf = slide.placeholders[1].text_frame
    tech = project_data.get('tech_stack_details', {})
    for k, v in tech.items():
        p = tf.add_paragraph()
        p.text = f"{k}: {v}"

    # Methodology
    if project_data.get('methodology'):
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "7. Methodology"
        slide.placeholders[1].text_frame.text = project_data.get('methodology', '')

    # Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "8. Conclusion"
    slide.placeholders[1].text_frame.text = "The system successfully implements the proposed architecture with advanced features and scalable design. Thank You!"

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
